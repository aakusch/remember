#!/usr/bin/env python3
"""
Generates the office-format fixtures for tests/parser-anydoc.test.ts.

Run with:
    python3 packages/core/tests/fixtures/office/generate-office-fixtures.py

Requires python-pptx and openpyxl; everything else is stdlib. The generated
binaries ARE committed — a .pptx/.xlsx/.odt/.epub is a zip archive, so byte
output is not reproducible across library versions (timestamps and entry
ordering vary). Committing them keeps the tests deterministic; this script
exists so the fixtures are auditable and regenerable rather than opaque blobs.

Deliberately NOT written into examples/multiformat-wiki/: that corpus is pinned
by `corpus_hash` in two committed benchmark artifacts
(benchmarks/results/remember-multiformat-wiki-*.json), and adding a file to it
would invalidate both.

Each document uses real named heading styles / real structure, because the
whole point of the fixture is to prove that container-declared heading levels
survive into `heading_path`.

No .doc fixture: the legacy OLE binary format cannot be synthesized without a
full Word writer. `.doc` support is covered by the extension-mapping test and
by the corrupt-input degradation test instead.
"""

import csv
import io
import os
import zipfile

HERE = os.path.dirname(os.path.abspath(__file__))


def out(name):
    return os.path.join(HERE, name)


# --------------------------------------------------------------------------
# .pptx — a deck. Slide titles are the heading signal.
# --------------------------------------------------------------------------
def make_pptx():
    from pptx import Presentation
    from pptx.util import Inches

    prs = Presentation()
    slides = [
        (
            "Incident response",
            [
                "Every production incident gets a severity within five minutes of detection.",
                "Sev1 pages the on-call engineer immediately.",
                "Sev2 waits for business hours unless customer data is at risk.",
            ],
        ),
        (
            "Escalation path",
            [
                "The on-call engineer owns the incident until it is handed over explicitly.",
                "Escalate to the platform lead after thirty minutes with no mitigation.",
                "Executive notification is required for any Sev1 lasting over two hours.",
            ],
        ),
        (
            "Postmortem",
            [
                "A written postmortem is due within five business days.",
                "Postmortems are blameless and name systems, not people.",
            ],
        ),
    ]

    layout = prs.slide_layouts[1]  # Title and Content
    for title, bullets in slides:
        slide = prs.slides.add_slide(layout)
        slide.shapes.title.text = title
        body = slide.placeholders[1].text_frame
        body.text = bullets[0]
        for line in bullets[1:]:
            body.add_paragraph().text = line

    # A slide carrying speaker notes, which anydoc extracts.
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = "Severity definitions"
    tb = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(3)).text_frame
    tb.text = "Sev1 means a full outage or confirmed data loss."
    tb.add_paragraph().text = "Sev3 means degraded performance with a workaround."
    slide.notes_slide.notes_text_frame.text = (
        "Remind the audience that severity is set by impact, never by effort."
    )

    prs.save(out("incident-deck.pptx"))


# --------------------------------------------------------------------------
# .xlsx — a sheet. Row-per-line is the expected serialization.
# --------------------------------------------------------------------------
def make_xlsx():
    from openpyxl import Workbook

    wb = Workbook()
    ws = wb.active
    ws.title = "Retention"
    rows = [
        ["Data class", "Retention period", "Owner"],
        ["Application logs", "30 days", "platform"],
        ["Audit logs", "7 years", "security"],
        ["Customer backups", "90 days", "platform"],
        ["Support transcripts", "2 years", "support"],
    ]
    for r in rows:
        ws.append(r)

    ws2 = wb.create_sheet("Contacts")
    ws2.append(["Team", "Escalation channel"])
    ws2.append(["platform", "#platform-oncall"])
    ws2.append(["security", "#security-incidents"])

    wb.save(out("retention-schedule.xlsx"))


# --------------------------------------------------------------------------
# ODF (.odt / .ods / .odp) — zip containers with a flat content.xml.
# --------------------------------------------------------------------------
ODF_STYLES = """<?xml version="1.0" encoding="UTF-8"?>
<office:document-styles
  xmlns:office="urn:oasis:names:tc:opendocument:xmlns:office:1.0"
  xmlns:style="urn:oasis:names:tc:opendocument:xmlns:style:1.0"
  xmlns:text="urn:oasis:names:tc:opendocument:xmlns:text:1.0"
  office:version="1.3"/>
"""


def write_odf(path, mimetype, content_xml):
    # The mimetype entry must be first and STORED (uncompressed) per the ODF
    # spec — that is the marker anydoc's container detection reads.
    with zipfile.ZipFile(path, "w", zipfile.ZIP_DEFLATED) as z:
        z.writestr(
            zipfile.ZipInfo("mimetype"), mimetype, compress_type=zipfile.ZIP_STORED
        )
        z.writestr(
            "META-INF/manifest.xml",
            '<?xml version="1.0" encoding="UTF-8"?>\n'
            '<manifest:manifest xmlns:manifest="urn:oasis:names:tc:opendocument:xmlns:manifest:1.0"'
            ' manifest:version="1.3">'
            f'<manifest:file-entry manifest:full-path="/" manifest:media-type="{mimetype}"/>'
            '<manifest:file-entry manifest:full-path="content.xml" manifest:media-type="text/xml"/>'
            "</manifest:manifest>",
        )
        z.writestr("styles.xml", ODF_STYLES)
        z.writestr("content.xml", content_xml)


ODF_NS = (
    ' xmlns:office="urn:oasis:names:tc:opendocument:xmlns:office:1.0"'
    ' xmlns:text="urn:oasis:names:tc:opendocument:xmlns:text:1.0"'
    ' xmlns:table="urn:oasis:names:tc:opendocument:xmlns:table:1.0"'
    ' xmlns:draw="urn:oasis:names:tc:opendocument:xmlns:drawing:1.0"'
    ' office:version="1.3"'
)


def make_odt():
    body = (
        '<text:h text:outline-level="1">Access control</text:h>'
        "<text:p>Access is granted by role, never to an individual account directly.</text:p>"
        '<text:h text:outline-level="2">Requesting access</text:h>'
        "<text:p>Open a request in the access portal and name the role you need.</text:p>"
        "<text:p>Your manager approves the request before security reviews it.</text:p>"
        '<text:h text:outline-level="3">Break-glass accounts</text:h>'
        "<text:p>Break-glass credentials live in the sealed vault and page security on use.</text:p>"
        '<text:h text:outline-level="2">Revocation</text:h>'
        "<text:p>Access is revoked automatically when a role assignment ends.</text:p>"
    )
    content = (
        '<?xml version="1.0" encoding="UTF-8"?>\n'
        f"<office:document-content{ODF_NS}>"
        f"<office:body><office:text>{body}</office:text></office:body>"
        "</office:document-content>"
    )
    write_odf(
        out("access-control.odt"),
        "application/vnd.oasis.opendocument.text",
        content,
    )


def make_ods():
    def row(cells):
        return "<table:table-row>" + "".join(
            f'<table:table-cell office:value-type="string"><text:p>{c}</text:p></table:table-cell>'
            for c in cells
        ) + "</table:table-row>"

    table = (
        '<table:table table:name="Budgets">'
        + row(["Team", "Quarterly budget", "Approver"])
        + row(["platform", "120000 USD", "cto"])
        + row(["support", "45000 USD", "coo"])
        + "</table:table>"
    )
    content = (
        '<?xml version="1.0" encoding="UTF-8"?>\n'
        f"<office:document-content{ODF_NS}>"
        f"<office:body><office:spreadsheet>{table}</office:spreadsheet></office:body>"
        "</office:document-content>"
    )
    write_odf(
        out("budgets.ods"),
        "application/vnd.oasis.opendocument.spreadsheet",
        content,
    )


def make_odp():
    def page(name, title, lines):
        frames = f'<draw:frame><draw:text-box><text:p>{title}</text:p></draw:text-box></draw:frame>'
        for line in lines:
            frames += f'<draw:frame><draw:text-box><text:p>{line}</text:p></draw:text-box></draw:frame>'
        return f'<draw:page draw:name="{name}">{frames}</draw:page>'

    pages = page(
        "page1",
        "Onboarding overview",
        [
            "Week one is environment setup and reading the architecture guide.",
            "Week two is a scoped starter task with a named reviewer.",
        ],
    ) + page(
        "page2",
        "Buddy system",
        ["Every new engineer is paired with a buddy for the first month."],
    )
    content = (
        '<?xml version="1.0" encoding="UTF-8"?>\n'
        f"<office:document-content{ODF_NS}>"
        f"<office:body><office:presentation>{pages}</office:presentation></office:body>"
        "</office:document-content>"
    )
    write_odf(
        out("onboarding-overview.odp"),
        "application/vnd.oasis.opendocument.presentation",
        content,
    )


# --------------------------------------------------------------------------
# .rtf — plain text container, headings via \outlinelevel.
# --------------------------------------------------------------------------
def make_rtf():
    doc = r"""{\rtf1\ansi\deff0
{\fonttbl{\f0 Calibri;}}
{\pard\s1\outlinelevel0\b\fs36 Secrets management\par}
{\pard\fs22 Secrets are stored in the managed vault and never in environment files committed to git.\par}
{\pard\s2\outlinelevel1\b\fs28 Rotation\par}
{\pard\fs22 Every credential rotates at least every ninety days.\par}
{\pard\fs22 A leaked credential is rotated immediately and the incident is logged.\par}
{\pard\s2\outlinelevel1\b\fs28 Exceptions\par}
{\pard\fs22 A rotation exception requires written security approval and an expiry date.\par}
}"""
    with open(out("secrets-management.rtf"), "w", encoding="ascii") as f:
        f.write(doc)


# --------------------------------------------------------------------------
# .epub — zip with a mimetype marker, an OPF spine, and XHTML chapters.
# --------------------------------------------------------------------------
def make_epub():
    chapter = """<?xml version="1.0" encoding="utf-8"?>
<html xmlns="http://www.w3.org/1999/xhtml"><head><title>Deploy guide</title></head>
<body>
<h1>Deploy guide</h1>
<p>Deploys run from the main branch only, through the release pipeline.</p>
<h2>Rollback</h2>
<p>Any deploy can be rolled back to the previous release from the pipeline UI.</p>
<p>A rollback is always safe: migrations are additive and backward compatible.</p>
<h2>Freeze windows</h2>
<p>No deploys land in the final week of a quarter without an approved exception.</p>
</body></html>
"""
    opf = """<?xml version="1.0" encoding="utf-8"?>
<package xmlns="http://www.idpf.org/2007/opf" version="3.0" unique-identifier="bookid">
  <metadata xmlns:dc="http://purl.org/dc/elements/1.1/">
    <dc:title>Deploy guide</dc:title>
    <dc:language>en</dc:language>
    <dc:identifier id="bookid">urn:uuid:remember-fixture-epub</dc:identifier>
  </metadata>
  <manifest>
    <item id="c1" href="chapter1.xhtml" media-type="application/xhtml+xml"/>
  </manifest>
  <spine><itemref idref="c1"/></spine>
</package>
"""
    container = """<?xml version="1.0" encoding="UTF-8"?>
<container version="1.0" xmlns="urn:oasis:names:tc:opendocument:xmlns:container">
  <rootfiles><rootfile full-path="OEBPS/content.opf" media-type="application/oebps-package+xml"/></rootfiles>
</container>
"""
    path = out("deploy-guide.epub")
    with zipfile.ZipFile(path, "w", zipfile.ZIP_DEFLATED) as z:
        z.writestr(
            zipfile.ZipInfo("mimetype"),
            "application/epub+zip",
            compress_type=zipfile.ZIP_STORED,
        )
        z.writestr("META-INF/container.xml", container)
        z.writestr("OEBPS/content.opf", opf)
        z.writestr("OEBPS/chapter1.xhtml", chapter)


# --------------------------------------------------------------------------
# .csv — no signature, so the parser must name the format explicitly.
# --------------------------------------------------------------------------
def make_csv():
    rows = [
        ["service", "tier", "oncall rotation"],
        ["checkout-api", "tier1", "platform-primary"],
        ["search-api", "tier1", "search-primary"],
        ["admin-console", "tier3", "none"],
        # A cell containing a comma and a quote, to prove escaping survives.
        ['billing-worker', 'tier2', 'finance, then platform'],
    ]
    buf = io.StringIO()
    csv.writer(buf, lineterminator="\n").writerows(rows)
    with open(out("service-catalog.csv"), "w", encoding="utf-8") as f:
        f.write(buf.getvalue())


def main():
    make_pptx()
    make_xlsx()
    make_odt()
    make_ods()
    make_odp()
    make_rtf()
    make_epub()
    make_csv()
    for name in sorted(os.listdir(HERE)):
        if name != os.path.basename(__file__):
            print(f"{name}\t{os.path.getsize(out(name))} bytes")


if __name__ == "__main__":
    main()
